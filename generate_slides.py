"""Generate PPTX slides for the GitHub Copilot workshop."""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Theme colors
BG_DARK = RGBColor(0x0F, 0x0F, 0x1E)
BG_SLIDE = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_MAIN = RGBColor(0xE0, 0xE0, 0xE0)
TEXT_HEADING = RGBColor(0xD4, 0xA5, 0x74)
TEXT_SUB = RGBColor(0xC4, 0x9A, 0x6C)
TEXT_DIM = RGBColor(0xB0, 0xB0, 0xB0)
TEXT_CODE = RGBColor(0xE8, 0xC4, 0x7C)
TEXT_ACCENT = RGBColor(0x7E, 0xB8, 0xDA)
BG_CODE = RGBColor(0x16, 0x16, 0x2A)
BG_TABLE_HDR = RGBColor(0x2D, 0x2D, 0x44)
BG_TABLE_ROW = RGBColor(0x1E, 0x1E, 0x36)

FONT_BODY = "Calibri"
FONT_CODE = "Consolas"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox.text_frame


def add_title(tf, text, size=Pt(36), color=TEXT_HEADING):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = FONT_BODY
    return p


def add_para(
    tf,
    text,
    size=Pt(18),
    color=TEXT_MAIN,
    bold=False,
    space_before=Pt(6),
    bullet=False,
    font_name=FONT_BODY,
):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    if bullet:
        p.level = 0
    return p


def add_code_block(tf, lines, size=Pt(14)):
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = TEXT_CODE
        p.font.name = FONT_CODE
        p.space_before = Pt(2)


def make_title_slide(prs, title, subtitle="", is_section=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    tf = add_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(2))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_title(tf, title, size=Pt(44) if not is_section else Pt(40))
    if subtitle:
        add_para(tf, subtitle, size=Pt(24), color=TEXT_SUB, space_before=Pt(16))
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER


def make_content_slide(prs, title, bullets, code=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    tf = add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
    add_title(tf, title, size=Pt(32))

    y = Inches(1.3)
    if bullets:
        tf2 = add_text_box(
            slide, Inches(0.6), y, Inches(8.8), Inches(3.5) if code else Inches(5.5)
        )
        first = True
        for b in bullets:
            if first:
                tf2.paragraphs[0].text = b
                tf2.paragraphs[0].font.size = Pt(18)
                tf2.paragraphs[0].font.color.rgb = TEXT_MAIN
                tf2.paragraphs[0].font.name = FONT_BODY
                tf2.paragraphs[0].space_before = Pt(4)
                first = False
            else:
                bold = b.startswith("**")
                text = b.strip("*") if bold else b
                add_para(tf2, text, bold=bold, space_before=Pt(8))
        y = Inches(4.8) if code else y

    if code:
        tf3 = add_text_box(slide, Inches(0.6), y, Inches(8.8), Inches(2.5))
        add_code_block(tf3, code)


def make_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    tf = add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
    add_title(tf, title, size=Pt(32))

    cols = len(headers)
    n_rows = len(rows) + 1
    col_width = Inches(8.8 / cols)
    table_shape = slide.shapes.add_table(
        n_rows, cols, Inches(0.6), Inches(1.4), Inches(8.8), Inches(0.45 * n_rows)
    )
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_TABLE_HDR
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_HEADING
            p.font.bold = True
            p.font.name = FONT_BODY

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_TABLE_ROW
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = TEXT_MAIN
                p.font.name = FONT_BODY


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- SLIDE 1: Title ---
    make_title_slide(
        prs,
        "GitHub Copilot Workshop",
        "From Beginner to Intermediate\n\nAn interactive, hands-on session with a real codebase\n\n~95 minutes  |  9 exercises  |  Live coding",
    )

    # --- SLIDE 2: Agenda ---
    make_table_slide(
        prs,
        "Agenda",
        ["Part", "Topics", "Time"],
        [
            ["1. Foundations", "Install, inline, Chat, Ask/Edit/Agent modes", "25 min"],
            ["2. Configuration", "Custom instructions, settings, model picker", "20 min"],
            ["3. Reusable workflows", "Agent mode, prompt files, custom agents", "25 min"],
            [
                "4. Advanced",
                "MCP, Copilot Coding Agent, context management",
                "15 min",
            ],
            ["5. Real-world", "End-to-end feature + PR", "5 min"],
            ["Wrap-up", "Quick wins, Q&A", "5 min"],
        ],
    )

    # --- SLIDE 3: Setup ---
    make_content_slide(
        prs,
        "Setup Check",
        [],
        code=[
            "# VS Code + GitHub Copilot + GitHub Copilot Chat extensions",
            "# Signed into GitHub with an active Copilot subscription",
            "",
            "python --version         # Python 3.11+",
            "uv --version             # uv package manager",
            "",
            "git clone https://github.com/zachschillaci27/github-copilot-workshop.git",
            "cd github-copilot-workshop",
            "uv sync",
            "uv run pytest            # 18 tests should pass",
            "code .                   # open in VS Code",
        ],
    )

    # --- SLIDE 4: Part 1 section ---
    make_title_slide(prs, "Part 1", "Foundations", is_section=True)

    # --- SLIDE 5: What is Copilot ---
    make_content_slide(
        prs,
        "Exercise 1 — Getting Started",
        [
            "What is GitHub Copilot?",
            "",
            "A family of AI features for software engineering",
            "Inline suggestions — ghost-text autocompletion as you type",
            "Chat — ask questions, explain code, explore",
            "Edits — multi-file diffs you approve hunk-by-hunk",
            "Agent — reads, edits, runs terminal, iterates until done",
            "Coding Agent — async, issue-driven, opens PRs for you",
        ],
    )

    # --- SLIDE 6: Surfaces ---
    make_table_slide(
        prs,
        "Where Copilot Lives",
        ["Surface", "What you use it for"],
        [
            ["VS Code / JetBrains / Visual Studio", "Inline + Chat + agent mode"],
            ["GitHub.com", "PR summaries, code review, issues"],
            ["Copilot Coding Agent", "Async: assign an issue, get a PR"],
            ["gh copilot / copilot CLI", "Terminal-native suggestions + chat"],
            ["GitHub Mobile", "On-the-go Q&A"],
        ],
    )

    # --- SLIDE 7: Ask / Edit / Agent ---
    make_table_slide(
        prs,
        "Ask / Edit / Agent — When to Use Which",
        ["Mode", "Does", "Use when"],
        [
            ["Ask", "Answers, no edits", "Exploring, explaining"],
            ["Edit", "Proposes multi-file diff you approve", "Focused refactor"],
            ["Agent", "Reads + edits + runs tools + loops", "Multi-step tasks"],
        ],
    )

    # --- SLIDE 8: Try It ---
    make_content_slide(
        prs,
        "Try It",
        [
            "Open the Chat view (Ctrl/Cmd+Alt+I), switch to Agent mode:",
        ],
        code=[
            '"What does this project do?"',
            "",
            '"Show me all the API endpoints"',
            "",
            '"Add a description field to the /health endpoint that',
            '  returns \'TaskFlow health check\', then run the tests"',
            "",
            "# Click the Checkpoint button to revert the whole turn.",
        ],
    )

    # --- SLIDE 9: Context shortcuts ---
    make_table_slide(
        prs,
        "Attach Context & Run Commands",
        ["Shortcut", "What it does"],
        [
            ["#file:<path>", "Attach a specific file"],
            ["#selection", "Attach the current editor selection"],
            ["#codebase <q>", "Semantic search across the repo"],
            ["#changes", "Working-tree diff"],
            ["#problems", "Items in the Problems panel"],
            ["@workspace <q>", "Project-wide Q&A specialist"],
            ["/explain /fix /tests /doc /new", "Built-in slash commands"],
            ["Ctrl/Cmd+I", "Inline chat on the current selection"],
        ],
    )

    # --- SLIDE 10: Part 2 section ---
    make_title_slide(prs, "Part 2", "Configuration", is_section=True)

    # --- SLIDE 11: Custom instructions ---
    make_content_slide(
        prs,
        "Exercise 2 — Custom Instructions",
        [
            "The single most impactful file for team-wide productivity",
            "",
            "Loaded automatically for every Chat / agent request. Covers:",
            "  - How to build and test your project",
            "  - Code conventions to follow",
            "  - Architecture overview",
            "  - Gotchas and workflows",
            "",
            "File locations:",
            "  .github/copilot-instructions.md       — Repo-wide (shared via git)",
            "  .github/instructions/*.instructions.md — Path-scoped (applyTo glob)",
            "  VS Code user settings                 — Personal, not shared",
        ],
    )

    # --- SLIDE 12: instructions content ---
    make_content_slide(
        prs,
        "What to Put in copilot-instructions.md",
        [],
        code=[
            "# TaskFlow API — Copilot Instructions",
            "",
            "## Build & Run",
            "- Install: `uv sync`",
            "- Run tests: `uv run pytest`",
            "- Lint: `uv run ruff check src/ tests/`",
            "",
            "## Code Conventions",
            "- Type hints on every function signature",
            "- Prefer `str | None` over `Optional[str]`",
            "- Keep endpoint handlers thin",
            "",
            "## API Design",
            "- All endpoints under /api/v1/",
            "- 201 for create, 204 for delete, 404 w/ detail for missing",
        ],
    )

    # --- SLIDE 13: path-scoped instructions ---
    make_content_slide(
        prs,
        "Path-Scoped Instructions",
        [
            "Only injected when the conversation touches matching files.",
        ],
        code=[
            "---",
            'applyTo: "tests/**/*.py"',
            'description: "Pytest conventions for TaskFlow"',
            "---",
            "",
            "# Test Conventions",
            "- Use the `client` fixture from tests/conftest.py",
            "- Name tests test_<action>_<scenario>",
            "- Every endpoint: happy path + error + one edge case",
        ],
    )

    # --- SLIDE 14: Settings ---
    make_content_slide(
        prs,
        "Exercise 3 — Settings & Permissions",
        [
            "VS Code settings shape Copilot's behaviour across modes.",
            "",
            ".vscode/settings.json                    — Workspace (shared via git)",
            "VS Code user settings                    — Personal, not shared",
            "GitHub org / enterprise console          — Policy, content exclusion",
            "",
            "Key Copilot settings:",
            "  github.copilot.enable                  — per-language toggle",
            "  chat.agent.enabled / maxRequests       — enable + cap tool calls",
            "  chat.tools.terminal.autoApprove        — allow/deny terminal cmds",
            "  chat.checkpoints.enabled               — revertable agent edits",
        ],
    )

    # --- SLIDE 15: Auto-approve rules ---
    make_content_slide(
        prs,
        "Terminal Auto-Approve",
        [
            "The safest way to reduce approval fatigue — allowlist the benign,",
            "denylist the dangerous, prompt on everything else.",
        ],
        code=[
            "{",
            '  "chat.tools.terminal.autoApprove": {',
            '    "uv run pytest": true,',
            '    "uv run ruff": true,',
            '    "uv sync": true,',
            '    "/^git (status|diff|log|branch|show)\\\\b/": true,',
            '    "rm": false,',
            '    "sudo": false,',
            '    "curl": false',
            "  }",
            "}",
        ],
    )

    # --- SLIDE 16: Part 3 section ---
    make_title_slide(prs, "Part 3", "Reusable workflows", is_section=True)

    # --- SLIDE 17: Agent mode deep-dive ---
    make_content_slide(
        prs,
        "Exercise 4 — Agent Mode",
        [
            "Copilot reads, plans, edits across files, runs the terminal, and loops.",
            "",
            "Safety rails:",
            "  Checkpoint button          — one-click revert of a whole turn",
            "  chat.agent.maxRequests     — cap on tool calls (default 25)",
            "  terminal / edit auto-approve — allow/deny lists",
            "  network filter (optional)  — HTTP domain allowlist",
            "",
            "No `PreToolUse` hooks like Claude Code — enforce via:",
            "  1. Instructions + agent judgement (soft guardrail)",
            "  2. Content exclusion policy (Business/Enterprise)",
            "  3. Pre-commit hooks (hard gate before commit)",
        ],
    )

    # --- SLIDE 18: Prompt files ---
    make_content_slide(
        prs,
        "Exercise 5 — Prompt Files",
        [
            "Reusable /slash commands, stored in .github/prompts/",
        ],
        code=[
            "---",
            "agent: agent",
            "description: Scaffold a new TaskFlow endpoint with tests",
            "---",
            "",
            "# Add New Endpoint",
            "",
            "Create a new API endpoint for: ${input:spec:e.g. GET /api/v1/tasks/overdue}",
            "",
            "1. Read src/taskflow/routers/tasks.py for the pattern",
            "2. Add the endpoint + database method",
            "3. Write tests using the `client` fixture",
            "4. Run `uv run pytest` to confirm green",
        ],
    )

    # --- SLIDE 19: Prompt frontmatter ---
    make_table_slide(
        prs,
        "Prompt File Frontmatter",
        ["Field", "Purpose"],
        [
            ["agent", "ask | agent | plan | <custom-agent-name>"],
            ["description", "Shown in the slash-command menu"],
            ["tools", "(Optional) allow-list of tools / MCP globs"],
            ["model", "(Optional) pin a specific model"],
            ["${input:name:placeholder}", "Prompt the user for input"],
            ["${selection} / ${file}", "Editor context injection"],
            ["#codebase / #changes / #file:…", "Chat-variable context injection"],
        ],
    )

    # --- SLIDE 20: Try prompts ---
    make_content_slide(
        prs,
        "Try the Pre-Built Prompts",
        [],
        code=[
            "/review src/taskflow/routers/tasks.py",
            "# -> Structured code review against project standards",
            "",
            "/test-coverage src/taskflow/routers/users.py",
            "# -> Finds untested functions, proposes test cases",
            "",
            "/add-endpoint GET /api/v1/tasks/overdue",
            "# -> Scaffolds endpoint + database method + tests",
            "",
            "/debug GET /api/v1/tasks returns 500 on invalid status",
            "# -> Investigation with root-cause + proposed fix",
            "",
            "/changelog",
            "# -> Keep-a-Changelog entry from #changes",
        ],
    )

    # --- SLIDE 21: Custom agents ---
    make_content_slide(
        prs,
        "Exercise 6 — Custom Agents",
        [
            "Persistent personas with scoped tools. Sit alongside built-in Ask",
            "and Agent in the agent picker. Renamed from 'chat modes' in 2026.",
            "",
            ".github/agents/<name>.agent.md",
        ],
        code=[
            "---",
            "description: Read-only planner — designs without editing",
            "tools: ['search/codebase', 'search/usages', 'web/fetch']",
            "---",
            "",
            "You are in planning mode. Produce a concrete implementation",
            "plan. Do not edit any files. Cite existing code with path:line.",
        ],
    )

    # --- SLIDE 22: Prompts vs modes ---
    make_table_slide(
        prs,
        "Prompt Files vs Custom Agents",
        ["Prompt file", "Custom agent"],
        [
            ["One-shot reusable task", "Persistent working persona"],
            ["Needs ${input:…} placeholders", "Multi-turn conversation"],
            ["Ships in /slash menu", "Ships in agent picker"],
            ["Invocation: /review src/…", "Invocation: switch agent, chat freely"],
            ["Good for: scaffolding, review, docs", "Good for: planner, reviewer, test-writer"],
        ],
    )

    # --- SLIDE 23: Part 4 section ---
    make_title_slide(prs, "Part 4", "Advanced", is_section=True)

    # --- SLIDE 24: MCP ---
    make_content_slide(
        prs,
        "Exercise 7 — MCP Servers",
        [
            "Extend agent mode with external tools (GitHub, databases, browsers).",
            "Cross-ecosystem: Copilot / Claude Code / Cursor all speak the same protocol.",
            "",
            "Configuration lives in .vscode/mcp.json (workspace) or user settings.",
            "Credentials come from `inputs` → ${input:id} — never hardcoded.",
        ],
        code=[
            "{",
            '  "inputs": [{ "id": "github-pat", "type": "promptString",',
            '               "password": true }],',
            '  "servers": {',
            '    "github": {',
            '      "type": "http",',
            '      "url": "https://api.githubcopilot.com/mcp",',
            '      "headers": {',
            '        "Authorization": "Bearer ${input:github-pat}"',
            "} } } }",
        ],
    )

    # --- SLIDE 25: Coding Agent ---
    make_content_slide(
        prs,
        "Copilot Coding Agent (Async)",
        [
            "Different product from agent mode — runs on GitHub Actions, not locally.",
            "",
            "Trigger: assign an issue to @copilot, or @copilot-mention in a comment.",
            "",
            "What happens:",
            "  1. Researches the repo, posts an implementation plan",
            "  2. Commits iteratively on a feature branch",
            "  3. Opens a draft PR, runs checks",
            "  4. Responds to @copilot feedback comments",
            "",
            "Environment bootstrap:",
            "  .github/workflows/copilot-setup-steps.yml (must be on default branch)",
        ],
    )

    # --- SLIDE 26: Context management ---
    make_content_slide(
        prs,
        "Exercise 8 — Context Management",
        [
            "Control what Copilot sees. Pick the right model for the turn.",
        ],
        code=[
            "# Chat variables",
            "#file:<path>           attach a file",
            "#selection             editor selection",
            "#codebase <q>          semantic repo search",
            "#changes               working-tree diff",
            "#problems              Problems panel",
            "#terminalLastCommand   last shell command + output",
            "",
            "# Participants",
            "@workspace   project-wide specialist",
            "@terminal    shell-aware",
            "@github      issues, PRs, repos",
        ],
    )

    # --- SLIDE 27: Model picker ---
    make_table_slide(
        prs,
        "Model Picker — Trade Latency for Depth",
        ["Tier", "Use for"],
        [
            ["Fast / general", "Default turn; quick edits; explanations"],
            ["Deep reasoning", "Plans, ambiguous refactors, architecture Q&A"],
            ["Lightweight", "High-volume small completions"],
            ["Vision", "Screenshots, diagrams, UI reviews"],
        ],
    )

    # --- SLIDE 28: Part 5 section ---
    make_title_slide(prs, "Part 5", "Real-world workflow", is_section=True)

    # --- SLIDE 29: Real-world workflow ---
    make_content_slide(
        prs,
        "Exercise 9 — End-to-End Feature",
        [
            "Putting it all together — local agent mode:",
            "",
            "1. Branch — create feature/task-search",
            "2. Planner mode — read-only plan with path:line citations",
            "3. /add-endpoint — scaffold route, database method, tests",
            "4. Reviewer mode — structured review against conventions",
            "5. /test-coverage — fill the gaps",
            "6. Commit & push, then @github open a PR",
            "",
            "Or go async: file the issue, assign to @copilot, review the PR.",
        ],
        code=[
            '"Add a search query param to GET /api/v1/tasks. Update the',
            " database layer, the route, and the tests. Use agent mode",
            ' end-to-end."',
        ],
    )

    # --- SLIDE 30: Wrap-up ---
    make_title_slide(prs, "Wrap-Up", is_section=True)

    # --- SLIDE 31: Quick wins ---
    make_content_slide(
        prs,
        "Quick Wins to Take Home",
        [
            "Today:",
            "  1. Add .github/copilot-instructions.md to your repos — 5 minutes, huge ROI",
            "",
            "This week:",
            "  2. Allowlist your common commands in chat.tools.terminal.autoApprove",
            "  3. Create one prompt file for your most repeated workflow",
            "",
            "This month:",
            "  4. Build a planner / reviewer custom agent for your team",
            "  5. Try assigning a small issue to @copilot and review the PR",
        ],
    )

    # --- SLIDE 32: Config landscape ---
    make_table_slide(
        prs,
        "The Copilot Configuration Landscape",
        ["Feature", "File", "Purpose"],
        [
            ["Repo instructions", ".github/copilot-instructions.md", "Team-wide conventions"],
            ["Path instructions", ".github/instructions/*.instructions.md", "Scoped rules (applyTo)"],
            ["Prompt files", ".github/prompts/*.prompt.md", "Reusable /slash commands"],
            ["Custom agents", ".github/agents/*.agent.md", "Scoped personas"],
            ["Workspace settings", ".vscode/settings.json", "Enablement + auto-approve"],
            ["MCP", ".vscode/mcp.json", "External tools"],
            ["Coding agent env", ".github/workflows/copilot-setup-steps.yml", "Cloud-agent bootstrap"],
        ],
    )

    # --- SLIDE 33: Interop ---
    make_content_slide(
        prs,
        "Interop With Other Agents",
        [
            "This repo also ships CLAUDE.md and .claude/skills/ — not used by",
            "Copilot, kept to show the same project instrumented for two tools.",
            "",
            "Keep conventions in sync across:",
            "  .github/copilot-instructions.md   (Copilot reads)",
            "  CLAUDE.md                         (Claude Code reads)",
            "  AGENTS.md                         (several tools recognize)",
            "",
            "MCP servers work across Copilot, Claude, Cursor, Windsurf — one config,",
            "many consumers.",
        ],
    )

    # --- SLIDE 34: Resources ---
    make_content_slide(
        prs,
        "Resources",
        [
            "GitHub Copilot docs: https://docs.github.com/en/copilot",
            "Copilot in VS Code:  https://code.visualstudio.com/docs/copilot/overview",
            "Community catalog:   https://github.com/github/awesome-copilot",
            "",
            "This repo:",
            "  exercises/ — 9 self-paced exercises",
            "  .github/   — all Copilot configuration",
            "  .vscode/   — workspace settings + MCP",
        ],
    )

    # --- SLIDE 35: Questions ---
    make_title_slide(
        prs,
        "Questions?",
        'Open Copilot Chat and ask:\n"@workspace what should I explore next?"',
    )

    prs.save("slides.pptx")
    print(f"Generated slides.pptx with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
